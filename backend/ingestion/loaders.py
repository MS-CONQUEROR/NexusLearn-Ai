import os
import pypdf
import docx
from pptx import Presentation
from langchain_core.documents import Document
from typing import List


class PDFLoader:
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        documents = []
        try:
            reader = pypdf.PdfReader(self.file_path)
        except Exception as e:
            raise ValueError(f"Could not read PDF: {e}")

        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if not text or not text.strip():
                continue
            documents.append(Document(
                page_content=text.strip(),
                metadata={
                    "source": os.path.basename(self.file_path),
                    "page": i + 1,
                    "source_type": "pdf"
                }
            ))
        return documents


class DOCXLoader:
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        documents = []
        try:
            doc = docx.Document(self.file_path)
        except Exception as e:
            raise ValueError(f"Could not read DOCX: {e}")

        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text.strip())

        # also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    full_text.append(row_text)

        if full_text:
            documents.append(Document(
                page_content="\n".join(full_text),
                metadata={
                    "source": os.path.basename(self.file_path),
                    "page": 1,
                    "source_type": "docx"
                }
            ))
        return documents


class PPTXLoader:
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        documents = []
        try:
            prs = Presentation(self.file_path)
        except Exception as e:
            raise ValueError(f"Could not read PPTX: {e}")

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                documents.append(Document(
                    page_content="\n".join(slide_text),
                    metadata={
                        "source": os.path.basename(self.file_path),
                        "page": i + 1,
                        "source_type": "pptx"
                    }
                ))
        return documents


def load_document(file_path: str) -> List[Document]:
    """
    Auto-detects file type and loads it.
    This is the single function the rest of the app calls.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return PDFLoader(file_path).load()
    elif ext == ".docx":
        return DOCXLoader(file_path).load()
    elif ext == ".pptx":
        return PPTXLoader(file_path).load()
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: pdf, docx, pptx")